import os
import shutil
import time
import comtypes.client as cli
import win32com.client
from pptx import Presentation
desktop = os.path.join(os.path.join(os.environ['USERPROFILE']), 'Desktop')


def rearrange_folders():
    """The function takes care of the old JUNC folder on the desktop: It checks if there is an old folder with some
    files -> if there is, it moves them to a new folder with the date stamp of creation date;
    If there is a folder but it's empty -> delete and create a fresh folder (important for saving the creation date)
    If there is no old JUNC folder -> it creates a new one.
    If there is a folder with the same timestamp, it creates another one, and adds an index number to it.
    """

    global desktop
    old_path = desktop + "\×JUNC×"
    if os.path.exists(old_path):
        old_id_path = old_path + "\×ID×.png"

        if len(os.listdir(old_path)) == 0:
            os.rmdir(old_path)
            os.makedirs(old_path)
        else:
            if os.path.exists(old_id_path):
                path_created = os.path.getmtime(old_id_path)
            else:
                path_created = os.path.getmtime(old_path)
            year, month, day, hour, minute, second = time.localtime(path_created)[:-3]
            folder = "\×JUNC× " + str(" %02d⌟%02d⌟%d %02d∶%02d" % (day, month, year, hour, minute))
            new_path = desktop + folder

            if not os.path.exists(new_path):
                os.makedirs(new_path)
            else:
                i = 1
                new_folder = False
                while not new_folder:
                    if not os.path.exists(new_path + "[%s]" % i):
                        new_path = new_path + "[%s]" % i
                        os.makedirs(new_path)
                        new_folder = True
                    else:
                        i += 1
            file_names = os.listdir(old_path)
            for file_name in file_names:
                shutil.move(os.path.join(old_path, file_name), new_path)
    else:
        os.makedirs(old_path)


def create_new_diagram_template_file():
    """The function creates a new copy of the Diagram templates; that will be the file to start working on
    It saves the new copy in the base directory"""
    src = os.getcwd() + r"\Diagram_template\Diagram_templates.pptx"
    prs = Presentation(src)
    prs.save("Diagram_new_template.pptx")


def del_slides(pres, chosen_type):
    """The function gets the diagram file and the chosen type of diagram and deletes all slides that don't match that
    number """
    i = len(pres.slides)
    while i > 0:
        if i != chosen_type:
            rId = pres.slides._sldIdLst[i - 1].rId
            pres.part.drop_rel(rId)
            del pres.slides._sldIdLst[i - 1]
        i -= 1
    pres.save("Del_Diagram.pptx")


def save_diagram(pres):
    """The function gets the final diagram pptx file and saves it in the created ×JUNC× folder on the desktop"""
    global desktop
    path = desktop + "\×JUNC×"
    save_export_path = path + r'\×Diagram×.pptx'
    pres.save(save_export_path)
    export_png(save_export_path)


def export_png(path):
    """The function creates a folder and exports a png photo of the diagram to that folder; before that,
    it calls 'fix_text' """
    f = os.path.abspath(path)
    fix_text(path)
    powerpoint = cli.CreateObject('Powerpoint.Application')
    powerpoint.ActivePresentation.Export(f, 'PNG')
    powerpoint.ActivePresentation.Close()


def fix_text(received):
    """The function fixes an issue occurs when trying to apply TEXT_TO_FIT_SHAPE in the other functions.
    It uses a different type of library (not python-pptx) that acts like VBA"""
    App = win32com.client.Dispatch("PowerPoint.Application")
    App.Visible = True
    Pres = App.Presentations.Open(received)
    for _ in Pres.Slides:
        App.CommandBars.ExecuteMso("SlideReset")
        App.CommandBars.ExecuteMso("SlideReset")
    Pres.Save()


def delete_temp_junc_pres():
    """The function deletes all the temporary powerpoint files created while creating the final one"""
    prs_to_del = ["Diagram_new_template.pptx", "Del_Diagram.pptx", "Street_Diagram.pptx", "Morn_Diagram.pptx",
                  "Eve_Diagram.pptx",
                  "Dirc_Diagram.pptx"]
    for temp_prs in prs_to_del:
        if os.path.exists(temp_prs):
            os.remove(temp_prs)
